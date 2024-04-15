import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from scipy.stats import mannwhitneyu
from termcolor import colored
from utils import load_all_developers_dataset


def process_input(input_text, uploaded_file, program_end_date=None, event_name=None):
    try:
        print(colored("Processing input...", "blue"))
        if uploaded_file is not None:
            print(colored("Reading from uploaded file...", "blue"))
            file_content = uploaded_file.read().decode("utf-8")
            github_handles = [
                handle.strip() for handle in file_content.split("\n") if handle.strip()
            ]
        else:
            github_handles = [handle.strip() for handle in input_text.split(",")]
            print(colored(f"GitHub handles: {github_handles}", "blue"))

        if program_end_date == "":
            program_end_date = None

        df = load_all_developers_dataset()

        print(colored("Filtering dataset...", "blue"))
        one_year_ago = pd.Timestamp.now() - pd.DateOffset(years=1)
        filtered_df = df[
            (df["developer"].isin(github_handles)) & (df["month_year"] >= one_year_ago)
        ]
        filtered_df = filtered_df.sort_values(by=["developer", "month_year"])
        filtered_df.loc[:, "month_year"] = pd.to_datetime(filtered_df["month_year"])
        line_fig = create_line_plot(filtered_df, github_handles, program_end_date)

        # Debug
        # print(colored("Debugging filtered dataset and github handles...", "blue"))
        # print(filtered_df.head(100))
        # print(filtered_df["developer"].unique())
        # print(github_handles)
        filtered_df.to_csv("debug.csv", index=False)
        # Debug

        analysis_result = perform_statistical_analysis(
            filtered_df, github_handles, program_end_date
        )

        new_developers_count = count_new_developers(
            filtered_df, github_handles, program_end_date
        )

        last_3_months = pd.Timestamp.now() - pd.DateOffset(months=3)
        recent_activity_user = filtered_df[filtered_df["month_year"] >= last_3_months]
        all_devs_df = load_all_developers_dataset()
        all_devs_filtered_df = all_devs_df[(all_devs_df["month_year"] >= last_3_months)]
        other_devs_recent_activity = all_devs_filtered_df[
            ~all_devs_filtered_df["developer"].isin(github_handles)
        ]

        user_specified_active = recent_activity_user[
            recent_activity_user["total_commits"] > 0
        ]
        other_developers_active = other_devs_recent_activity[
            other_devs_recent_activity["total_commits"] > 0
        ]

        box_fig = create_box_plot(user_specified_active, other_developers_active)

        print(colored("Classifying developers...", "blue"))
        classification_df = classify_developers(github_handles, recent_activity_user)
        print(colored("Classification completed.", "blue"))

        comparison_result = compare_user_developers_to_others(
            user_specified_active, other_developers_active, df, program_end_date
        )

        growth_rate_result = compare_growth_rate(
            user_specified_active, other_developers_active, df
        )

        tldr_summary = generate_tldr_summary(
            github_handles,
            classification_df,
            analysis_result,
            new_developers_count,
            comparison_result,
            growth_rate_result,
            event_name,
        )

        return (
            line_fig,
            box_fig,
            classification_df,
            analysis_result,
            new_developers_count,
            comparison_result,
            growth_rate_result,
            tldr_summary,
        )
    except Exception as e:
        print(colored(f"Error processing input: {e}", "red"))
        return (
            None,
            None,
            None,
            None,
            "Error in processing input. Check logs for more details on the error",
            None,
            None,
            "Error in processing input. Check logs for more details on the error",
        )


def create_line_plot(filtered_df, github_handles, program_end_date):
    plot_df = filtered_df.copy()
    missing_developers = set(github_handles) - set(plot_df["developer"].unique())
    for developer in missing_developers:
        new_row = pd.DataFrame(
            {
                "developer": [developer],
                "month_year": [pd.Timestamp.now()],
                "total_commits": [0],
            }
        )
        plot_df = pd.concat([plot_df, new_row], ignore_index=True)
    plot_df = (
        plot_df.groupby(["developer", "month_year"])["total_commits"]
        .sum()
        .reset_index()
    )
    line_fig = px.line(
        plot_df,
        x="month_year",
        y="total_commits",
        color="developer",
        labels={"month_year": "Month", "total_commits": "Number of Commits"},
        title="Commits per Month",
    )
    if program_end_date:
        program_end_date = pd.to_datetime(program_end_date)
        line_fig.add_vline(
            x=program_end_date, line_width=2, line_dash="dash", line_color="red"
        )
    return line_fig


def create_box_plot(user_specified_active, other_developers_active):
    box_fig = go.Figure()
    box_fig.add_trace(
        go.Box(
            y=user_specified_active["total_commits"], name="User Specified Developers"
        )
    )
    box_fig.add_trace(
        go.Box(y=other_developers_active["total_commits"], name="Other Developers")
    )
    box_fig.update_layout(
        title="Comparison of Monthly Commits in the Last 3 Months: User Specified vs. Other Developers (Active Only)",
        yaxis_title="Total Monthly Commits",
        yaxis=dict(range=[0, 50]),
    )
    return box_fig


def classify_developers(github_handles, recent_activity_user):
    classification = []
    for handle in github_handles:
        dev_df = recent_activity_user[recent_activity_user["developer"] == handle]
        total_recent_commits = dev_df["total_commits"].sum()
        if dev_df.empty or total_recent_commits == 0:
            status = "Always been inactive"
        elif total_recent_commits < 10:
            status = "Low-level active (<10 commits)"
        elif total_recent_commits < 20:
            status = "Moderately active (10-19 commits)"
        else:
            status = "Highly involved (20+ commits)"
        classification.append((handle, status, total_recent_commits))

    sort_keys = {
        "Highly involved": 1,
        "Moderately active": 2,
        "Low-level active": 3,
        "Previously active but no longer": 4,
        "Always been inactive": 5,
    }
    classification_df = pd.DataFrame(
        classification, columns=["Developer", "Classification", "Total Recent Commits"]
    )
    classification_df["Sort Key"] = classification_df["Classification"].map(sort_keys)
    classification_df.sort_values(
        by=["Sort Key", "Total Recent Commits"], ascending=[True, False], inplace=True
    )
    classification_df.drop(["Sort Key", "Total Recent Commits"], axis=1, inplace=True)
    return classification_df


def perform_statistical_analysis(filtered_df, github_handles, program_end_date_str):
    if program_end_date_str is None:
        return "Program end date not provided. Unable to perform statistical analysis."

    program_end_date = pd.to_datetime(program_end_date_str)
    before_program = filtered_df[filtered_df["month_year"] < program_end_date]
    after_program = filtered_df[filtered_df["month_year"] >= program_end_date]

    before_counts = before_program.groupby("developer")["total_commits"].median()
    after_counts = after_program.groupby("developer")["total_commits"].median()

    all_developers = pd.Series(0, index=github_handles)
    before_counts = before_counts.reindex(all_developers.index, fill_value=0)
    after_counts = after_counts.reindex(all_developers.index, fill_value=0)

    if (before_counts == 0).all() or (after_counts == 0).all():
        return "Not enough data for statistical analysis. All values are zero in either before or after counts."

    stat, p_value = mannwhitneyu(after_counts, before_counts)
    analysis_result = (
        f"Mann-Whitney U test statistic: {stat:.3f}, P-value: {p_value:.3f}\n"
    )

    if p_value < 0.2:
        if stat > 0:
            analysis_result += (
                "CONGRATULATIONS! The commit activity from the builders taking part in this program is statistically higher after the program."
                "Difference in commit activity before and after the program is considered significant."
            )
        else:
            analysis_result += (
                "UNFORTUNATELY: The commit activity from the builders taking part in this program is lower after the program."
                "Difference in commit activity before and after the program is considered significant."
            )
    else:
        analysis_result += "NOTHING ESPECIAL. No significant difference in commit activity before and after the program."

    return analysis_result


def count_new_developers(filtered_df, github_handles, program_end_date_str):
    if program_end_date_str is None:
        print(
            colored(
                "Program end date not provided. Unable to count new developers. No problem.",
                "yellow",
            )
        )
        return (
            "Program end date not provided. Unable to count new developers. No problem."
        )

    program_end_date = pd.to_datetime(program_end_date_str)
    two_months_after_program = program_end_date + pd.DateOffset(months=2)

    before_program = filtered_df[filtered_df["month_year"] < program_end_date]
    after_program = filtered_df[
        (filtered_df["month_year"] >= program_end_date)
        & (filtered_df["month_year"] <= two_months_after_program)
    ]

    before_developers = before_program["developer"].unique()
    after_developers = after_program["developer"].unique()

    new_developers = set(after_developers) - set(before_developers)
    new_developers_str = ", ".join(new_developers)

    return f"Number of new developers committing code within 2 months after the program: {len(new_developers)}\nNew developers: {new_developers_str}"


def compare_user_developers_to_others(
    user_specified_active, other_developers_active, df, program_end_date_str
):
    if program_end_date_str is None:
        print(
            colored(
                "Program end date not provided. Unable to compare user-specified developers to others. No problem.",
                "yellow",
            )
        )
        return "Program end date not provided. Unable to compare user-specified developers to others. No problem."

    program_end_date = pd.to_datetime(program_end_date_str)
    user_commits = df[
        (df["developer"].isin(user_specified_active["developer"]))
        & (df["month_year"] >= program_end_date)
    ]["total_commits"]
    other_commits = df[
        (df["developer"].isin(other_developers_active["developer"]))
        & (df["month_year"] >= program_end_date)
    ]["total_commits"]

    if len(user_commits) == 0 or len(other_commits) == 0:
        print(
            colored(
                "Not enough data for comparison. Either user-specified developers or developers in the database have no commits after the program end date. Update database",
                "red",
            )
        )

    stat, p_value = mannwhitneyu(user_commits, other_commits)
    comparison_result = (
        f"Mann-Whitney U test statistic: {stat:.3f}, P-value: {p_value:.3f}\n"
    )

    if p_value < 0.25:
        if stat > 0:
            comparison_result += "GREAT! The user-specified developers have a significantly higher number of commits compared to other developers since the program end date."
        else:
            comparison_result += "BAD! The user-specified developers have a significantly lower number of commits compared to other developers since the program end date."
    else:
        comparison_result += "NOTHING SPECIAL! There is no significant difference in the number of commits between user-specified developers and other developers since the program end date."

    return comparison_result


def compare_growth_rate(user_specified_active, other_developers_active, df):
    user_growth_rates = []
    other_growth_rates = []

    for developer in user_specified_active["developer"].unique():
        user_df = df[df["developer"] == developer]
        user_df = user_df.sort_values("month_year")
        user_commits = user_df["total_commits"].tolist()
        user_growth_rate = calculate_average_growth_rate(user_commits)
        user_growth_rates.append(user_growth_rate)

    for developer in other_developers_active["developer"].unique():
        other_df = df[df["developer"] == developer]
        other_df = other_df.sort_values("month_year")
        other_commits = other_df["total_commits"].tolist()
        other_growth_rate = calculate_average_growth_rate(other_commits)
        other_growth_rates.append(other_growth_rate)

    stat, p_value = mannwhitneyu(user_growth_rates, other_growth_rates)
    comparison_result = (
        f"Mann-Whitney U test statistic: {stat:.3f}, P-value: {p_value:.3f}\n"
    )

    if p_value < 0.25:
        if stat > 0:
            comparison_result += "GOOD! These developers have a significantly higher average growth rate of commit activity compared to other developers."
        else:
            comparison_result += "BAD! These developers have a significantly lower average growth rate of commit activity compared to other developers."
    else:
        comparison_result += "NOTHING SPECIAL! There is no significant difference in the average growth rate of commit activity between user-specified developers and other developers."

    return comparison_result


def calculate_average_growth_rate(commits):
    growth_rates = []
    for i in range(1, len(commits)):
        if commits[i - 1] != 0:
            growth_rate = (commits[i] - commits[i - 1]) / commits[i - 1]
            growth_rates.append(growth_rate)
    if len(growth_rates) > 0:
        return sum(growth_rates) / len(growth_rates)
    else:
        return 0


def generate_tldr_summary(
    github_handles,
    classification_df,
    analysis_result,
    new_developers_count,
    comparison_result,
    growth_rate_result,
    event_name,
):
    summary = f"### TLDR Summary for {', '.join(github_handles)}\n\n"
    highly_involved_devs = classification_df[
        classification_df["Classification"] == "Highly involved"
    ]["Developer"].tolist()
    if highly_involved_devs:
        summary += f"**High Performers:** {', '.join(highly_involved_devs)}\n\n"
    if "higher after the program" in analysis_result:
        summary += "**Commit Activity from partipants in the program:** GREAT! Significantly higher after the program.\n\n"
    elif "lower after the program" in analysis_result:
        summary += "**Commit Activity from partipants in the program:** BAD! Significantly lower after the program.\n\n"
    else:
        summary += "**Commit Activity from partipants in the program:** NOTHING SPECIAL: No significant change after the program.\n\n"
    if new_developers_count.startswith("Number of new developers"):
        summary += (
            f"**New Developers:** {new_developers_count.split(':')[1].strip()}\n\n"
        )
    if "significantly higher number of commits" in comparison_result:
        summary += "**Comparison with Other Developers:** GREAT! User-specified developers have a significantly higher number of commits.\n\n"
    elif "significantly lower number of commits" in comparison_result:
        summary += "**Comparison with Other Developers:** BAD! User-specified developers have a significantly lower number of commits.\n\n"
    else:
        summary += "**Comparison with Other Developers:** NOTHING SPECIAL: No significant difference in the number of commits.\n\n"
    if "significantly higher average growth rate" in growth_rate_result:
        summary += "**Growth Rate in number of commits:** GREAT! User-specified developers have a significantly higher average growth rate.\n\n"
    elif "significantly lower average growth rate" in growth_rate_result:
        summary += "**Growth Rate in number of commits:** BAD! User-specified developers have a significantly lower average growth rate.\n\n"
    else:
        summary += "**Growth Rate in number of commits:** NOTHING SPECIAL: No significant difference in the average growth rate.\n\n"
    if event_name:
        summary += f"*Note: The analysis is based on the {event_name} event.*\n\n"
    return summary


def create_ppt_report(
    tldr_summary,
    line_fig,
    box_fig,
    classification_df,
    analysis_result,
    new_developers_count,
    comparison_result,
    growth_rate_result,
):
    prs = Presentation()

    # Add TLDR summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "TLDR Summary"
    content = slide.placeholders[1]
    content.text = tldr_summary

    # Add line plot slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Commits per Month"
    line_fig.write_image("line_plot.png")
    slide.shapes.add_picture(
        "line_plot.png", Inches(1), Inches(2), width=Inches(8), height=Inches(5)
    )

    # Add box plot slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Box Plot Comparison (Last 3 Months)"
    box_fig.write_image("box_plot.png")
    slide.shapes.add_picture(
        "box_plot.png", Inches(1), Inches(2), width=Inches(8), height=Inches(5)
    )

    # Add developer classification slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Which developers are the most active?"
    content = slide.placeholders[1]
    content.text = "\n".join(
        [
            f"{row['Developer']}: {row['Classification']}"
            for _, row in classification_df.iterrows()
        ]
    )

    # Add statistical analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Are these developers committing more code after the program than before the program?"
    content = slide.placeholders[1]
    content.text = analysis_result

    # Add new developers count slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Do we have new developers after the program?"
    content = slide.placeholders[1]
    content.text = new_developers_count

    # Add comparison with other developers slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Do the developers of this program commit more code than other Starknet developers?"
    content = slide.placeholders[1]
    content.text = comparison_result

    # Add growth rate comparison slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Is the increase rate in commits from these developers higher than that of other Starknet developers?"
    content = slide.placeholders[1]
    content.text = growth_rate_result

    prs.save("developer_insights_report.pptx")


def main():
    df = load_all_developers_dataset()
    max_available_month = df["month_year"].max().strftime("%Y-%m")

    try:
        with open("./github-metrics/assets/style.css") as f:
            st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
    except:
        try:
            with open("../assets/style.css") as f:
                st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
        except:
            with open("assets/style.css") as f:
                st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

    st.title("Starknet Star Tracker: GitHub Starknet Developer Insights")
    st.markdown(
        """
    üìù **Note:** This tool is developed and maintained by Omar Espejel (Telegram: @espejelomar) from the Starknet Foundation. Feel free to contact him to provide feedback or ask questions.
    """
    )
    st.markdown(
        """
        This tool allows you to analyze the GitHub activity of developers within the Starknet ecosystem.
        Enter GitHub handles separated by commas or upload a CSV file with GitHub handles in a single column
        to see their monthly commit activity, involvement classification, and comparisons with other develope        """
    )
    st.markdown(
        """
        üì∫ **Video Tutorial:** Please watch this [5-minute video tutorial](https://www.loom.com/share/b60e7f1bd1ee473b97e9c84c74df692a) examining an African Bootcamp and the Basecamp bootcamp as examples to start using the app effectively.
        """
    )

    text_input = st.text_input(
        "Enter GitHub handles separated by commas",
        placeholder="e.g., user1,user2,user3",
    )
    file_input = st.file_uploader(
        "Or upload a CSV file with GitHub handles in a single column",
        type=["csv"],
    )
    st.markdown(
        """
        *Note:* When uploading a CSV, ensure it contains a single column of GitHub handles without a header row.
        """
    )

    program_end_date_input = st.text_input(
        f"Program End Date (YYYY-MM-DD) (Max available: {max_available_month})",
        placeholder="e.g., 2023-06-30",
    )
    event_name_input = st.text_input(
        "Event Name (optional)",
        placeholder="e.g., Basecamp, Hackathon",
    )
    st.markdown(
        """
        üí° *Tip: Specifying a program end date allows you to analyze the impact of events like Basecamp or Hackathons on developer activity. Leave it blank to analyze overall activity.*
        """
    )

    if st.button("Analyze"):
        if program_end_date_input > max_available_month:
            st.warning(
                f"The specified date {program_end_date_input} is not available in the dataset. Please choose a date up to {max_available_month}."
            )
        else:
            (
                line_fig,
                box_fig,
                classification_df,
                analysis_result,
                new_developers_count,
                comparison_result,
                growth_rate_result,
                tldr_summary,
            ) = process_input(
                text_input, file_input, program_end_date_input, event_name_input
            )

        st.markdown(tldr_summary)

        st.plotly_chart(line_fig)
        st.plotly_chart(box_fig)

        with st.expander("üèÜ Which developers are the most active?"):
            st.dataframe(classification_df)
            st.markdown(
                """
                ### Developer Classification Criteria
                - **Always been inactive**: No commits have been recorded in the dataset.
                - **Previously active but no longer**: Had commits earlier but none in the last 3 months.
                - **Low-level active**: Fewer than 10 commits in the last 3 months.
                - **Moderately active**: 10 to 19 commits in the last 3 months.
                - **Highly involved**: 20 or more commits in the last 3 months.
                """
            )

        with st.expander("üÜï Do we have new developers after the program?"):
            st.text(new_developers_count)

        with st.expander(
            "üìä Are these developers committing more code after the program than before the program?"
        ):
            st.text(analysis_result)
            st.markdown(
                """
                The Mann-Whitney U test is used to compare the commit activity of developers before and after the program.
                - The test statistic measures the difference in the distribution of commits between the two groups (before and after).
                - The p-value indicates the probability of observing such a difference by chance, assuming there is no real difference between the groups.
                - A p-value less than 0.2 suggests that the difference is considered significant.
                - A positive test statistic indicates that the commit activity is higher after the program, while a negative value indicates lower activity.
                """
            )

        with st.expander(
            "üîç Do the developers of this program commit more code than other Starknet developers?"
        ):
            st.text(comparison_result)
            st.markdown(
                """
                The Mann-Whitney U test is used to compare the commit activity of the user-specified developers with the rest of the developers in the database since the program end date.
                - The test statistic measures the difference in the distribution of commits between the two groups.
                - The p-value indicates the probability of observing such a difference by chance, assuming there is no real difference between the groups.
                - A p-value less than 0.25 suggests that the difference is considered significant.
                - If the test statistic is positive, it means the user-specified developers have a higher number of commits compared to other developers, and vice versa.
                """
            )

        with st.expander(
            "üìà Is the increase rate in commits from these developers higher than that of other Starknet developers?"
        ):
            st.text(growth_rate_result)
            st.markdown(
                """
                The average growth rate of commit activity is compared between the user-specified developers and other developers.
                - The growth rate is calculated as the relative change in the number of commits from one month to the next.
                - The Mann-Whitney U test is used to compare the average growth rates between the two groups.
                - A p-value less than 0.25 suggests that the difference in average growth rates is statistically significant.
                - If the test statistic is positive, it means the user-specified developers have a higher average growth rate compared to other developers, and vice versa.
                """
            )

        st.markdown(
            """
            üí° *Disclaimer: This information is only for open-source repos and should be taken with a grain of salt. Commits in certain repos may be more important than others, and there are many private repos from several teams that are not included in this analysis.*
            """
        )

        create_ppt_report(
            tldr_summary,
            line_fig,
            box_fig,
            classification_df,
            analysis_result,
            new_developers_count,
            comparison_result,
            growth_rate_result,
        )
        st.success("Report exported as developer_insights_report.pptx")


if __name__ == "__main__":
    main()
